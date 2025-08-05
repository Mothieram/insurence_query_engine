import os
import logging
import json
import re
from typing import List, Dict, Optional, Tuple
import pdfplumber
import docx2txt
from pptx import Presentation
from email import policy
from email.parser import BytesParser
from datetime import datetime
from langchain.text_splitter import RecursiveCharacterTextSplitter
import spacy
from spacy.matcher import PhraseMatcher


logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('insurance_clause_detector.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)


# Load NLP model (medium model for better performance)
try:
    nlp = spacy.load("en_core_web_md")
except OSError:
    logger.error("Spacy model 'en_core_web_md' not found. Please install it first.")
    raise

# Configure logging


class InsuranceClauseDetector:
    """Advanced insurance clause detection with NLP-enhanced parsing."""
    
    def __init__(self):
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=150,
            separators=self._get_insurance_separators()
        )
        
        # Initialize phrase matcher for insurance-specific terms
        self.matcher = PhraseMatcher(nlp.vocab, attr="LOWER")
        self._initialize_patterns()
        
    def _get_insurance_separators(self) -> List[str]:
        """Insurance-specific text boundaries for chunking."""
        return [
            "\n## ", "\n** ", "\nSECTION ", "\nCLAUSE ",  # Section headers
            "\n1. ", "\n2. ", "\n3. ", "\n(a) ", "\n(b) ",  # Numbered clauses
            "\n• ", "\n- ", "\n* ",  # Bullet points
            "\n\n", "\n", " "  # Standard separators
        ]
    
    def _initialize_patterns(self):
        """Initialize NLP patterns for insurance clause detection."""
        # Coverage triggers
        coverage_terms = ["covered", "benefit", "eligible", "payable", "reimbursement",
                         "coverage includes", "we will pay", "policy provides"]
        self._add_phrases_to_matcher("COVERAGE", coverage_terms)
        
        # Exclusion triggers
        exclusion_terms = ["not covered", "excluded", "exception", "not payable",
                          "does not apply", "will not pay", "not eligible"]
        self._add_phrases_to_matcher("EXCLUSION", exclusion_terms)
        
        # Limitation triggers
        limitation_terms = ["limit", "maximum", "cap", "up to", "not exceeding",
                          "subject to", "restricted to", "only if"]
        self._add_phrases_to_matcher("LIMITATION", limitation_terms)
        
        # Condition triggers
        condition_terms = ["provided that", "if", "when", "subject to", "contingent upon",
                         "depending on", "requires", "must", "shall"]
        self._add_phrases_to_matcher("CONDITION", condition_terms)
        
        # Definition triggers
        definition_terms = ["means", "refers to", "defined as", "for purposes of",
                          "in this policy", "shall mean"]
        self._add_phrases_to_matcher("DEFINITION", definition_terms)
    
    def _add_phrases_to_matcher(self, label: str, phrases: List[str]):
        """Helper to add phrases to the NLP matcher."""
        patterns = [nlp(text) for text in phrases]
        self.matcher.add(label, patterns)
    
    def process_document(self, file_path: str) -> Optional[Dict]:
        """Main method to process insurance documents and detect clauses."""
        try:
            logger.info(f"Processing document: {file_path}")
            
            if not os.path.exists(file_path):
                logger.error(f"File not found: {file_path}")
                return None
                
            raw_text = self._extract_raw_text(file_path)
            if not raw_text:
                logger.warning(f"No extractable text found in: {file_path}")
                return None
                
            # Enhanced processing
            doc_info = {
                "metadata": self._extract_metadata(file_path, raw_text),
                "clauses": self._detect_clauses(raw_text),
                "analysis": self._analyze_document(raw_text)
            }
            
            return doc_info
            
        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}", exc_info=True)
            return None

    def _detect_clauses(self, text: str) -> List[Dict]:
        """Detect and classify insurance clauses with NLP."""
        chunks = self.splitter.split_text(text)
        clauses = []
        
        for chunk in chunks:
            doc = nlp(chunk)
            matches = self.matcher(doc)
            
            # Get unique clause types found in this chunk
            found_types = set()
            for match_id, start, end in matches:
                found_types.add(nlp.vocab.strings[match_id])
            
            # Classify the overall clause
            clause_type = self._determine_primary_clause_type(found_types, chunk)
            
            # Extract important entities
            entities = self._extract_insurance_entities(doc)
            
            clauses.append({
                "text": chunk.strip(),
                "type": clause_type,
                "subtypes": list(found_types),
                "entities": entities,
                "confidence": self._calculate_confidence(chunk, clause_type)
            })
            
        return clauses
    
    def _determine_primary_clause_type(self, found_types: set, text: str) -> str:
        """Determine the primary clause type based on multiple factors."""
        if not found_types:
            return self._fallback_classify_clause(text)
        
        # Prioritize certain types when multiple matches exist
        priority_order = ["EXCLUSION", "LIMITATION", "CONDITION", "COVERAGE", "DEFINITION"]
        
        for clause_type in priority_order:
            if clause_type in found_types:
                return clause_type
                
        return found_types.pop()  # Return first if none in priority order
    
    def _fallback_classify_clause(self, text: str) -> str:
        """Fallback classification when no patterns are matched."""
        text_lower = text.lower()
        if any(x in text_lower for x in ['not cover', 'exclu', 'exception']):
            return "EXCLUSION"
        elif any(x in text_lower for x in ['cover', 'benefit', 'eligible']):
            return "COVERAGE"
        elif any(x in text_lower for x in ['limit', 'maximum', 'cap']):
            return "LIMITATION"
        elif any(x in text_lower for x in ['if', 'when', 'provided that']):
            return "CONDITION"
        elif any(x in text_lower for x in ['means', 'defined as']):
            return "DEFINITION"
        return "GENERAL"
    
    def _extract_insurance_entities(self, doc) -> List[Dict]:
        """Extract insurance-specific entities using NLP."""
        entities = []
        
        for ent in doc.ents:
            entities.append({
                "text": ent.text,
                "type": ent.label_,
                "start": ent.start_char,
                "end": ent.end_char
            })
        
        # Additional insurance-specific entity extraction
        monetary_terms = [token.text for token in doc if token.like_num and token.is_currency]
        if monetary_terms:
            entities.append({
                "text": " ".join(monetary_terms),
                "type": "MONETARY_AMOUNT",
                "start": 0,
                "end": 0  # Will be calculated in post-processing
            })
            
        return entities
    
    def _calculate_confidence(self, text: str, clause_type: str) -> float:
        """Calculate confidence score for clause classification."""
        # Simple implementation - can be enhanced with ML model
        keywords = {
            "COVERAGE": ["cover", "benefit", "pay"],
            "EXCLUSION": ["not cover", "excluded", "exception"],
            "LIMITATION": ["limit", "maximum", "cap"],
            "CONDITION": ["if", "when", "provided that"],
            "DEFINITION": ["means", "defined as"]
        }
        
        text_lower = text.lower()
        matches = sum(1 for word in keywords.get(clause_type, []) 
                    if word in text_lower)
        total_keywords = len(keywords.get(clause_type, []))
        
        return matches / total_keywords if total_keywords > 0 else 0.5
    
    def _analyze_document(self, text: str) -> Dict:
        """Perform document-level analysis of insurance clauses."""
        doc = nlp(text)
        sentences = list(doc.sents)
        
        return {
            "clause_distribution": self._get_clause_distribution(text),
            "coverage_analysis": self._analyze_coverage(text),
            "risk_keywords": self._extract_risk_keywords(text),
            "document_stats": {
                "sentences": len(sentences),
                "avg_sentence_length": sum(len(s.text) for s in sentences)/len(sentences) if sentences else 0,
                "readability": self._calculate_readability(text)
            }
        }
    
    def _get_clause_distribution(self, text: str) -> Dict:
        """Calculate distribution of clause types in the document."""
        clauses = self._detect_clauses(text)
        type_counts = {}
        
        for clause in clauses:
            clause_type = clause["type"]
            type_counts[clause_type] = type_counts.get(clause_type, 0) + 1
        
        return type_counts
    
    def _analyze_coverage(self, text: str) -> Dict:
        """Analyze coverage-related clauses."""
        coverage_clauses = [c for c in self._detect_clauses(text) 
                          if c["type"] in ["COVERAGE", "EXCLUSION", "LIMITATION"]]
        
        coverage_items = []
        for clause in coverage_clauses:
            doc = nlp(clause["text"])
            coverage_items.extend(self._extract_coverage_items(doc))
            
        return {
            "coverage_items": coverage_items,
            "total_coverages": len([c for c in coverage_clauses if c["type"] == "COVERAGE"]),
            "total_exclusions": len([c for c in coverage_clauses if c["type"] == "EXCLUSION"]),
            "total_limitations": len([c for c in coverage_clauses if c["type"] == "LIMITATION"])
        }
    
    def _extract_coverage_items(self, doc) -> List[str]:
        """Extract specific items covered/excluded from text."""
        items = []
        for chunk in doc.noun_chunks:
            if any(token.text.lower() in ["cover", "include", "exclude"] 
                  for token in chunk.root.head.children):
                items.append(chunk.text)
        return items
    
    def _extract_risk_keywords(self, text: str) -> List[str]:
        """Extract risk-related keywords from text."""
        risk_terms = [
            "risk", "hazard", "danger", "liability", "exposure", 
            "accident", "injury", "damage", "loss", "claim"
        ]
        return list(set(
            term for term in risk_terms 
            if re.search(rf'\b{term}\b', text, re.IGNORECASE)
        ))
    
    def _calculate_readability(self, text: str) -> float:
        """Simple readability score (Flesch-Kincaid approximation)."""
        sentences = text.split('.')
        words = text.split()
        
        if not sentences or not words:
            return 0
            
        avg_sentence_length = len(words) / len(sentences)
        avg_word_length = sum(len(word) for word in words) / len(words)
        
        # Simplified formula
        return max(0, min(100, 206.835 - 1.015 * avg_sentence_length - 84.6 * avg_word_length))
    
    # (Keep all the file extraction methods from the original code)
    def _extract_raw_text(self, file_path: str) -> Optional[str]:
        """Extract text from various file formats with robust error handling."""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            
            if ext == '.pdf':
                return self._extract_pdf_text(file_path)
            elif ext == '.docx':
                return docx2txt.process(file_path)
            elif ext == '.pptx':
                return self._extract_ppt_text(file_path)
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            elif ext == '.eml':
                return self._extract_email_text(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_path}")
                return None
                
        except Exception as e:
            logger.error(f"Text extraction failed for {file_path}: {str(e)}")
            return None

    def _extract_pdf_text(self, file_path: str) -> Optional[str]:
        """Specialized PDF text extraction with image-based fallback detection."""
        try:
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if not page_text:
                            if len(page.images) > 0:
                                logger.warning(f"Page {i+1} appears to be image-based")
                            else:
                                logger.warning(f"Page {i+1} returned no text")
                            continue
                        text_parts.append(page_text)
                    except Exception as page_error:
                        logger.error(f"Error processing page {i+1}: {str(page_error)}")
            return "\n".join(text_parts) if text_parts else None
        except pdfplumber.PDFSyntaxError:
            logger.error(f"Invalid PDF file: {file_path}")
            return None

    def _extract_ppt_text(self, file_path: str) -> str:
        """Extract text from PowerPoint files."""
        ppt = Presentation(file_path)
        text = []
        for i, slide in enumerate(ppt.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text.append(f"Slide {i+1}: " + " | ".join(slide_text))
        return "\n".join(text)

    def _extract_email_text(self, file_path: str) -> str:
        """Extract text content from email files."""
        with open(file_path, "rb") as f:
            msg = BytesParser(policy=policy.default).parse(f)
            body = []
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        charset = part.get_content_charset() or "utf-8"
                        body.append(part.get_payload(decode=True).decode(charset))
            else:
                charset = msg.get_content_charset() or "utf-8"
                body.append(msg.get_payload(decode=True).decode(charset))
            return "\n".join(body)

    def _extract_metadata(self, file_path: str, text: str) -> Dict:
        """Extract metadata from insurance documents."""
        metadata = {
            "source_file": os.path.basename(file_path),
            "policy_number": self._find_pattern(r'Policy\s*(No|Number)[:.]?\s*(\w+)', text),
            "effective_date": self._find_pattern(r'Effective\s*Date[:.]?\s*(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})', text),
            "coverage_types": self._detect_coverage_types(text),
            "file_info": {
                "path": file_path,
                "last_modified": str(datetime.fromtimestamp(os.path.getmtime(file_path)))
            }
        }
        return metadata

    def _find_pattern(self, pattern: str, text: str) -> Optional[str]:
        """Helper to extract text using regex patterns."""
        match = re.search(pattern, text, re.IGNORECASE)
        return match.group(2) if match else None

    def _detect_coverage_types(self, text: str) -> List[str]:
        """Identify insurance coverage types from text."""
        coverage_map = {
            "health": r'\b(medical|hospitalization|surgery|treatment)\b',
            "life": r'\b(life insurance|death benefit|mortality)\b',
            "accident": r'\b(accident|injury|disability)\b',
            "property": r'\b(property|fire|theft)\b',
            "liability": r'\b(liability|negligence|lawsuit)\b',
            "auto": r'\b(auto|vehicle|car|collision)\b'
        }
        return [cov_type for cov_type, pattern in coverage_map.items() 
                if re.search(pattern, text, re.IGNORECASE)]

def process_insurance_files(file_paths: List[str], output_dir: str) -> Dict:
    """Process batch of insurance documents with comprehensive reporting."""
    detector = InsuranceClauseDetector()
    os.makedirs(output_dir, exist_ok=True)
    
    results = {
        "success": [],
        "failed": [],
        "summary": {
            "total_files": len(file_paths),
            "successful": 0,
            "failed": 0,
            "clause_stats": {}
        }
    }
    
    for file_path in file_paths:
        file_result = {
            "file": file_path,
            "status": "pending",
            "output_path": None,
            "error": None
        }
        
        try:
            document_data = detector.process_document(file_path)
            if document_data:
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_path = os.path.join(output_dir, f"{base_name}_clauses.json")
                
                with open(output_path, 'w', encoding='utf-8') as f:
                    json.dump(document_data, f, indent=2)
                
                # Update clause statistics
                if "clauses" in document_data:
                    for clause in document_data["clauses"]:
                        clause_type = clause["type"]
                        results["summary"]["clause_stats"][clause_type] = \
                            results["summary"]["clause_stats"].get(clause_type, 0) + 1
                
                file_result.update({
                    "status": "success",
                    "output_path": output_path,
                    "metadata": document_data.get("metadata", {})
                })
                results["success"].append(file_result)
                results["summary"]["successful"] += 1
            else:
                file_result.update({
                    "status": "failed",
                    "error": "No data extracted"
                })
                results["failed"].append(file_result)
                results["summary"]["failed"] += 1
                
        except Exception as e:
            file_result.update({
                "status": "error",
                "error": str(e)
            })
            results["failed"].append(file_result)
            results["summary"]["failed"] += 1
            
    return results

def main():
    """Main execution with proper directory handling and reporting."""
    # Configure paths
    SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
    INPUT_DIR = os.path.abspath(os.path.join(SCRIPT_DIR, "..", "file_inputs"))
    OUTPUT_DIR = os.path.abspath(os.path.join(SCRIPT_DIR, "..", "processed_clauses"))
    
    # Verify input directory
    if not os.path.exists(INPUT_DIR):
        logger.error(f"Input directory not found: {INPUT_DIR}")
        logger.info(f"Current working directory: {os.getcwd()}")
        return
    
    # Gather input files (case-insensitive)
    input_files = []
    valid_extensions = ('.pdf', '.docx', '.pptx', '.txt', '.eml')
    for f in os.listdir(INPUT_DIR):
        if f.lower().endswith(valid_extensions):
            input_files.append(os.path.join(INPUT_DIR, f))
    
    if not input_files:
        logger.error(f"No supported files found in {INPUT_DIR}")
        logger.info("Directory contents:")
        for f in os.listdir(INPUT_DIR):
            logger.info(f" - {f}")
        return
    
    logger.info(f"\nStarting processing of {len(input_files)} insurance documents")
    start_time = datetime.now()
    
    # Process files
    results = process_insurance_files(input_files, OUTPUT_DIR)
    
    # Generate report
    duration = (datetime.now() - start_time).total_seconds()
    logger.info(f"\nProcessing completed in {duration:.2f} seconds")
    logger.info(f"Successful: {results['summary']['successful']}")
    logger.info(f"Failed: {results['summary']['failed']}")
    
    # Print clause statistics
    logger.info("\nClause Statistics:")
    for clause_type, count in results["summary"]["clause_stats"].items():
        logger.info(f"{clause_type}: {count}")
    
    if results['failed']:
        logger.info("\nFailed files:")
        for f in results['failed']:
            logger.info(f" - {f['file']}: {f.get('error', 'Unknown error')}")
    
    # Save full report
    report_path = os.path.join(OUTPUT_DIR, "clause_analysis_report.json")
    with open(report_path, 'w') as f:
        json.dump(results, f, indent=2)
    logger.info(f"\nFull report saved to: {report_path}")

if __name__ == "__main__":
    main()